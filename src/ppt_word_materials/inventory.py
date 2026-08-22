from __future__ import annotations

import hashlib
import io
import json
import re
from collections import defaultdict
from pathlib import Path

from pypdf import PdfReader
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .ppt_render import crop_shape_from_slide
from .ppt_render import render_pptx_slides
from .contracts import validate_deliverable_asset
from .ppt_objects import (
    bbox_xywh,
    bbox_xyxy,
    expand_bbox_with_nearby_text,
    shape_text_recursive,
    valid_bbox_xywh,
)
from .pdf_objects import (
    area_ratio as pdf_area_ratio,
    discover_pdf_object_regions,
    expand_rect_with_nearby_text,
    rect_bbox,
)
from .quality import evaluate_asset_quality


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sanitize_text(text: str) -> str:
    return re.sub(r"[\ud800-\udfff]", "\ufffd", text)


def _safe_stem(path: Path) -> str:
    return re.sub(r"[^0-9A-Za-z\u4e00-\u9fff_-]+", "-", path.stem).strip("-")


def _image_extension(blob: bytes) -> str | None:
    if blob.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if blob.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if blob.startswith((b"GIF87a", b"GIF89a")):
        return "gif"
    if blob.startswith(b"BM"):
        return "bmp"
    if blob.startswith((b"II*\x00", b"MM\x00*")):
        return "tiff"
    if len(blob) >= 12 and blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        return "webp"
    return None


def _verified_raster(blob: bytes) -> tuple[str | None, str | None]:
    extension = _image_extension(blob)
    if extension is None:
        return None, "unsupported_or_unknown_image_magic"
    try:
        with Image.open(io.BytesIO(blob)) as image:
            width, height = image.size
            if width <= 0 or height <= 0 or width * height > 100_000_000:
                return None, "unsafe_image_dimensions"
            image.verify()
    except (OSError, ValueError) as exc:
        return None, f"invalid_raster_image:{exc}"
    return extension, None


def _visual_hash(path: str | Path | None) -> str | None:
    if not path:
        return None
    try:
        with Image.open(path) as image:
            gray = image.convert("L").resize((9, 8), Image.Resampling.LANCZOS)
            pixels = list(gray.getdata())
    except (OSError, ValueError):
        return None
    bits = []
    for row in range(8):
        offset = row * 9
        bits.extend(
            pixels[offset + column] > pixels[offset + column + 1]
            for column in range(8)
        )
    value = sum(int(bit) << index for index, bit in enumerate(bits))
    return f"{value:016x}"


ENTITY_PATTERN = re.compile(
    r"[\u4e00-\u9fffA-Za-z]{2,8}?(?:资本|基金|公会|协会|公司|集团|研究院|研究所|科技局)"
)
NEARBY_ENTITY_PATTERN = re.compile(
    r"[\u4e00-\u9fffA-Za-z]{2,8}?(?:资本|公会|协会|集团|研究院|研究所|科技局)"
)
GENERIC_ENTITY_TERMS = {
    "社会资本", "政府投资基金", "政府引导基金", "私募投资基金", "投资基金",
    "母基金", "子基金", "产业基金", "并购基金", "基金", "公司", "子公司",
}
GENERIC_ENTITY_FRAGMENTS = (
    "投资基金", "引导基金", "产业基金", "并购基金", "母基金", "子基金",
    "社会资本", "证券基金", "股权基金",
)


def _valid_identity_entity(entity: str) -> bool:
    if entity in GENERIC_ENTITY_TERMS or entity.startswith(("年", "过", "与", "的")):
        return False
    if any(fragment in entity for fragment in GENERIC_ENTITY_FRAGMENTS):
        return False
    if entity.endswith("资本") and len(entity) > 6:
        return False
    return True


def _identity_entities(record: dict[str, object]) -> list[str]:
    values = [
        (str(record.get("object_text") or ""), ENTITY_PATTERN),
        (Path(str(record.get("source_file") or "")).stem, ENTITY_PATTERN),
        (str(record.get("nearby_text") or "")[:500], NEARBY_ENTITY_PATTERN),
    ]
    entities: list[str] = []
    for value, pattern in values:
        for entity in pattern.findall(value):
            if _valid_identity_entity(entity) and entity not in entities:
                entities.append(entity)
    return entities[:12]


def _corner_placement(record: dict[str, object]) -> bool:
    if "corner_placement" in record:
        return bool(record["corner_placement"])
    bbox = record.get("source_bbox")
    canvas = record.get("source_canvas")
    if not isinstance(bbox, list) or len(bbox) != 4:
        return False
    if not isinstance(canvas, list) or len(canvas) != 2:
        return False
    x0, y0, x1, y1 = map(float, bbox)
    width, height = map(float, canvas)
    if width <= 0 or height <= 0:
        return False
    horizontal_edge = x0 / width <= 0.16 or x1 / width >= 0.84
    vertical_edge = y0 / height <= 0.22 or y1 / height >= 0.78
    return horizontal_edge and vertical_edge


def _image_dimensions(record: dict[str, object]) -> tuple[int, int]:
    value = record.get("original_media_path") or record.get("render_path")
    if not value:
        return 0, 0
    try:
        with Image.open(str(value)) as image:
            return image.size
    except (OSError, ValueError):
        return 0, 0


def _assign_asset_role(record: dict[str, object]) -> None:
    scope = str(record.get("asset_scope") or "")
    asset_type = str(record.get("asset_type") or "")
    if scope == "context":
        record["asset_role"] = "context"
        record["role_reason"] = ["context_only"]
        record["role_confidence"] = 1.0
        return
    if scope == "rejected" or asset_type == "ppt_picture_unreadable":
        record["asset_role"] = "forbidden"
        record["role_reason"] = [str(record.get("extraction_error") or "unreadable")]
        record["role_confidence"] = 1.0
        return
    if asset_type not in {"ppt_picture", "pdf_image_crop"}:
        record["asset_role"] = "content_visual"
        record["role_reason"] = ["structured_visual_object"]
        record["role_confidence"] = 0.95
        return

    area_ratio = float(record.get("area_ratio") or 0)
    reuse_count = int(record.get("reuse_count") or 1)
    corner = _corner_placement(record)
    record["corner_placement"] = corner
    width, height = _image_dimensions(record)
    compact_enough = 0.002 <= area_ratio <= 0.03
    size_ok = not width or not height or (max(width, height) >= 130 and min(width, height) >= 40)
    shape_ok = not width or not height or 0.6 <= width / max(1, height) <= 8
    entities = _identity_entities(record)
    record["identity_entities"] = entities
    repeated_identity = reuse_count >= 3 and corner and compact_enough
    supported_reuse_two = reuse_count == 2 and corner and compact_enough and bool(entities)
    title_identity = (
        int(record.get("source_page_or_slide") or 0) == 1
        and corner
        and compact_enough
        and bool(entities)
    )
    if size_ok and shape_ok and (repeated_identity or supported_reuse_two or title_identity):
        record["asset_role"] = "identity_candidate"
        record["role_reason"] = [
            "reused_corner_identity" if reuse_count >= 3 else "entity_supported_corner_identity"
        ]
        record["role_confidence"] = 0.9 if reuse_count >= 3 else 0.75
        record["logo_family"] = entities[0] if entities else str(record.get("source_media_sha256") or record.get("sha256") or record.get("asset_id"))
        record["is_decorative"] = False
    elif area_ratio < 0.02:
        record["asset_role"] = "decoration"
        record["role_reason"] = ["small_unverified_picture"]
        record["role_confidence"] = 0.85
        record["is_decorative"] = True
    else:
        record["asset_role"] = "content_visual"
        record["role_reason"] = ["substantive_picture"]
        record["role_confidence"] = 0.8
        record["is_decorative"] = False


def deduplicate_assets(records: list[dict[str, object]]) -> list[dict[str, object]]:
    groups: dict[str, list[dict[str, object]]] = defaultdict(list)
    for record in records:
        digest = str(record.get("source_media_sha256") or record.get("sha256") or "")
        if digest:
            groups[digest].append(record)

    for group in groups.values():
        canonical = str(group[0]["asset_id"])
        reuse_count = len(group)
        for index, record in enumerate(group):
            record["reuse_count"] = reuse_count
            record["duplicate_of"] = None if index == 0 else canonical
    for record in records:
        record.setdefault("reuse_count", 1)
        record.setdefault("duplicate_of", None)
        record.setdefault("is_decorative", False)
        _assign_asset_role(record)
    return records


def extract_pdf_pages(
    source_pdf: str | Path, output_dir: str | Path, *, render: bool = True
) -> list[dict[str, object]]:
    source = Path(source_pdf)
    destination = Path(output_dir)
    destination.mkdir(parents=True, exist_ok=True)
    source_hash = _sha256_bytes(source.read_bytes())
    reader = PdfReader(source)
    import fitz

    pdf = fitz.open(source)
    render_dir = destination / "pages"
    crop_dir = destination / "object-crops"
    original_dir = destination / "object-assets"
    if render:
        render_dir.mkdir(parents=True, exist_ok=True)
        crop_dir.mkdir(parents=True, exist_ok=True)
    original_dir.mkdir(parents=True, exist_ok=True)

    records: list[dict[str, object]] = []
    for index, reader_page in enumerate(reader.pages, start=1):
        page = pdf[index - 1]
        text = sanitize_text((reader_page.extract_text() or "").strip())
        render_path = None
        if render:
            render_path = render_dir / f"page-{index:03d}.png"
            page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False).save(render_path)
        digest = (
            _sha256_bytes(render_path.read_bytes())
            if render_path is not None
            else _sha256_bytes(f"{source_hash}:pdf:{index}".encode("utf-8"))
        )
        records.append(
            {
                "asset_id": f"pdf-{source_hash[:12]}-p{index:03d}",
                "source_file": str(source.resolve()),
                "source_kind": "pdf",
                "source_page_or_slide": index,
                "asset_type": "pdf_page",
                "asset_scope": "context",
                "deliverable": False,
                "text": text,
                "nearby_text": text,
                "render_path": str(render_path.resolve()) if render_path else None,
                "original_media_path": None,
                "sha256": digest,
                "delivery_render_sha256": digest,
                "source_media_sha256": None,
                "area_ratio": 1.0,
                "is_decorative": False,
                "is_searchable": bool(text),
                "ocr_required": not bool(text),
            }
        )

        for object_index, region in enumerate(
            discover_pdf_object_regions(page), start=1
        ):
            object_rect = fitz.Rect(region["object_rect"])
            delivery_rect = expand_rect_with_nearby_text(page, object_rect)
            object_text = sanitize_text(page.get_text("text", clip=object_rect).strip())
            crop_text = sanitize_text(page.get_text("text", clip=delivery_rect).strip())
            asset_type = str(region["asset_type"])
            crop_path = None
            crop_digest = _sha256_bytes(
                f"{source_hash}:pdf:{index}:{asset_type}:{rect_bbox(delivery_rect)}".encode(
                    "utf-8"
                )
            )
            if render:
                crop_path = crop_dir / (
                    f"pdf-{source_hash[:12]}-p{index:03d}-{asset_type}-{object_index:03d}.png"
                )
                page.get_pixmap(
                    matrix=fitz.Matrix(2, 2), clip=delivery_rect, alpha=False
                ).save(crop_path)
                crop_digest = _sha256_bytes(crop_path.read_bytes())

            original_path = None
            source_media_sha256 = None
            xref = int(region.get("xref") or 0)
            if asset_type == "pdf_image_crop" and xref > 0:
                try:
                    extracted = pdf.extract_image(xref)
                    source_media_sha256 = _sha256_bytes(extracted["image"])
                    extension = str(extracted.get("ext") or "bin")
                    original_path = original_dir / f"xref-{xref}.{extension}"
                    if not original_path.exists():
                        original_path.write_bytes(extracted["image"])
                except (KeyError, RuntimeError, ValueError):
                    original_path = None

            object_ratio = pdf_area_ratio(object_rect, page.rect)
            records.append(
                {
                    "asset_id": f"pdf-{source_hash[:12]}-p{index:03d}-{asset_type}-{object_index:03d}",
                    "source_file": str(source.resolve()),
                    "source_kind": "pdf",
                    "source_page_or_slide": index,
                    "asset_type": asset_type,
                    "asset_scope": "visual_object",
                    "deliverable": True,
                    "text": crop_text or text,
                    "object_text": object_text,
                    "nearby_text": text,
                    "render_path": str(crop_path.resolve()) if crop_path else None,
                    "original_media_path": (
                        str(original_path.resolve()) if original_path else None
                    ),
                    "sha256": crop_digest,
                    "delivery_render_sha256": crop_digest,
                    "source_media_sha256": source_media_sha256,
                    "visual_phash": _visual_hash(original_path or crop_path),
                    "area_ratio": object_ratio,
                    "is_decorative": object_ratio < 0.02,
                    "is_searchable": bool(crop_text or text),
                    "ocr_required": not bool(crop_text or text),
                    "source_bbox": rect_bbox(delivery_rect),
                    "object_bbox": rect_bbox(object_rect),
                    "source_canvas": [float(page.rect.width), float(page.rect.height)],
                    "xref": xref or None,
                }
            )
    pdf.close()
    return deduplicate_assets(records)


def _shape_text(shape) -> str:
    return sanitize_text(shape_text_recursive(shape).strip())


def safe_picture_blob(shape) -> tuple[bytes | None, str | None, str | None]:
    first_error: str | None = None
    try:
        image = shape.image
        blob = image.blob
        extension, error = _verified_raster(blob)
        if error is None:
            return blob, extension, None
        first_error = error
    except (AttributeError, KeyError, ValueError) as exc:
        first_error = str(exc)

    try:
        blip = shape._pic.blipFill.blip
        relationship_id = getattr(blip, "embed", None)
        if not relationship_id and hasattr(blip, "get"):
            relationship_id = blip.get(qn("r:embed"))
        if not relationship_id:
            linked_id = getattr(blip, "link", None)
            if not linked_id and hasattr(blip, "get"):
                linked_id = blip.get(qn("r:link"))
            if linked_id:
                return None, None, "linked_image_not_embedded"
            return None, None, first_error or "missing_blip_relationship"
        related_part = shape.part.related_part(relationship_id)
        blob = bytes(related_part.blob)
        extension, error = _verified_raster(blob)
        if error is not None:
            return None, None, error
        return blob, extension, None
    except (AttributeError, KeyError, TypeError, ValueError) as exc:
        return None, None, first_error or str(exc)


def extract_pptx_structure(
    source_pptx: str | Path, output_dir: str | Path
) -> list[dict[str, object]]:
    source = Path(source_pptx)
    destination = Path(output_dir)
    media_dir = destination / "object-assets"
    media_dir.mkdir(parents=True, exist_ok=True)
    source_hash = _sha256_bytes(source.read_bytes())
    deck = Presentation(source)
    slide_area = float(deck.slide_width * deck.slide_height)
    records: list[dict[str, object]] = []

    for slide_number, slide in enumerate(deck.slides, start=1):
        texts = [text for shape in slide.shapes if (text := _shape_text(shape))]
        slide_text = "\n".join(texts)
        records.append(
            {
                "asset_id": f"ppt-{source_hash[:12]}-s{slide_number:03d}",
                "source_file": str(source.resolve()),
                "source_kind": "pptx",
                "source_page_or_slide": slide_number,
                "asset_type": "ppt_slide",
                "asset_scope": "context",
                "deliverable": False,
                "text": slide_text,
                "nearby_text": slide_text,
                "render_path": None,
                "original_media_path": None,
                "sha256": _sha256_bytes(
                    f"{source_hash}:ppt:{slide_number}".encode("utf-8")
                ),
                "delivery_render_sha256": None,
                "source_media_sha256": None,
                "area_ratio": 1.0,
                "is_decorative": False,
                "is_searchable": bool(slide_text),
                "ocr_required": not bool(slide_text),
            }
        )

        for shape_index, shape in enumerate(slide.shapes, start=1):
            area_ratio = float(shape.width * shape.height) / slide_area
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob, extension, extraction_error = safe_picture_blob(shape)
                if blob is None:
                    records.append(
                        {
                            "asset_id": f"ppt-{source_hash[:12]}-s{slide_number:03d}-unreadable{shape_index:03d}",
                            "source_file": str(source.resolve()),
                            "source_kind": "pptx",
                            "source_page_or_slide": slide_number,
                            "asset_type": "ppt_picture_unreadable",
                            "asset_scope": "rejected",
                            "deliverable": False,
                            "text": slide_text,
                            "nearby_text": slide_text,
                            "render_path": None,
                            "original_media_path": None,
                            "sha256": "",
                            "delivery_render_sha256": None,
                            "source_media_sha256": None,
                            "area_ratio": area_ratio,
                            "is_decorative": True,
                            "is_searchable": False,
                            "ocr_required": False,
                            "bbox_emu": [shape.left, shape.top, shape.width, shape.height],
                            "source_bbox": [
                                shape.left,
                                shape.top,
                                shape.left + shape.width,
                                shape.top + shape.height,
                            ],
                            "source_canvas": [deck.slide_width, deck.slide_height],
                            "extraction_error": extraction_error,
                        }
                    )
                    continue
                digest = _sha256_bytes(blob)
                target = media_dir / f"{digest[:16]}.{extension}"
                if not target.exists():
                    target.write_bytes(blob)
                records.append(
                    {
                        "asset_id": f"ppt-{source_hash[:12]}-s{slide_number:03d}-pic{shape_index:03d}",
                        "source_file": str(source.resolve()),
                        "source_kind": "pptx",
                        "source_page_or_slide": slide_number,
                        "asset_type": "ppt_picture",
                        "asset_scope": "visual_object",
                        "deliverable": True,
                        "text": slide_text,
                        "object_text": "",
                        "nearby_text": slide_text,
                        "render_path": str(target.resolve()),
                        "original_media_path": str(target.resolve()),
                        "sha256": digest,
                        "delivery_render_sha256": digest,
                        "source_media_sha256": digest,
                        "visual_phash": _visual_hash(target),
                        "area_ratio": area_ratio,
                        "is_decorative": area_ratio < 0.02,
                        "is_searchable": bool(slide_text),
                        "ocr_required": False,
                        "bbox_emu": [shape.left, shape.top, shape.width, shape.height],
                        "source_bbox": [
                            shape.left,
                            shape.top,
                            shape.left + shape.width,
                            shape.top + shape.height,
                        ],
                        "source_canvas": [deck.slide_width, deck.slide_height],
                    }
                )
            elif getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
                asset_type = "ppt_chart_crop" if shape.has_chart else "ppt_table_crop"
                object_bbox = bbox_xyxy(shape)
                delivery_bbox = expand_bbox_with_nearby_text(
                    shape,
                    slide.shapes,
                    slide_width=deck.slide_width,
                    slide_height=deck.slide_height,
                )
                records.append(
                    {
                        "asset_id": f"ppt-{source_hash[:12]}-s{slide_number:03d}-{asset_type}-{shape_index:03d}",
                        "source_file": str(source.resolve()),
                        "source_kind": "pptx",
                        "source_page_or_slide": slide_number,
                        "asset_type": asset_type,
                        "asset_scope": "visual_object",
                        "deliverable": True,
                        "text": _shape_text(shape) or slide_text,
                        "object_text": _shape_text(shape),
                        "nearby_text": slide_text,
                        "render_path": None,
                        "original_media_path": None,
                        "sha256": "",
                        "delivery_render_sha256": None,
                        "source_media_sha256": None,
                        "area_ratio": area_ratio,
                        "is_decorative": area_ratio < 0.02,
                        "is_searchable": bool(slide_text),
                        "ocr_required": False,
                        "bbox_emu": bbox_xywh(delivery_bbox),
                        "object_bbox_emu": bbox_xywh(object_bbox),
                        "source_bbox": delivery_bbox,
                        "source_canvas": [deck.slide_width, deck.slide_height],
                    }
                )
            elif shape.shape_type in {
                MSO_SHAPE_TYPE.GROUP,
                MSO_SHAPE_TYPE.DIAGRAM,
                MSO_SHAPE_TYPE.IGX_GRAPHIC,
                MSO_SHAPE_TYPE.CANVAS,
            }:
                object_bbox = bbox_xyxy(shape)
                delivery_bbox = expand_bbox_with_nearby_text(
                    shape,
                    slide.shapes,
                    slide_width=deck.slide_width,
                    slide_height=deck.slide_height,
                )
                records.append(
                    {
                        "asset_id": f"ppt-{source_hash[:12]}-s{slide_number:03d}-ppt_diagram_crop-{shape_index:03d}",
                        "source_file": str(source.resolve()),
                        "source_kind": "pptx",
                        "source_page_or_slide": slide_number,
                        "asset_type": "ppt_diagram_crop",
                        "asset_scope": "visual_object",
                        "deliverable": True,
                        "text": _shape_text(shape) or slide_text,
                        "object_text": _shape_text(shape),
                        "nearby_text": slide_text,
                        "render_path": None,
                        "original_media_path": None,
                        "sha256": "",
                        "delivery_render_sha256": None,
                        "source_media_sha256": None,
                        "area_ratio": area_ratio,
                        "is_decorative": area_ratio < 0.02,
                        "is_searchable": bool(_shape_text(shape) or slide_text),
                        "ocr_required": False,
                        "bbox_emu": bbox_xywh(delivery_bbox),
                        "object_bbox_emu": bbox_xywh(object_bbox),
                        "source_bbox": delivery_bbox,
                        "source_canvas": [deck.slide_width, deck.slide_height],
                    }
                )
    for record in records:
        bbox = record.get("bbox_emu")
        if (
            record.get("asset_scope") == "visual_object"
            and isinstance(bbox, list)
            and not valid_bbox_xywh(bbox)
        ):
            record["asset_scope"] = "rejected"
            record["deliverable"] = False
            record["render_path"] = None
            record["extraction_error"] = "invalid_object_geometry"
            record["is_decorative"] = True

    return deduplicate_assets(records)


def attach_pptx_renders(
    records: list[dict[str, object]],
    rendered_slides: list[str | Path],
    *,
    slide_size_emu: list[int | float],
    output_dir: str | Path,
) -> list[dict[str, object]]:
    slide_paths = {
        index: Path(path).resolve()
        for index, path in enumerate(rendered_slides, start=1)
    }
    crop_dir = Path(output_dir)
    crop_dir.mkdir(parents=True, exist_ok=True)

    for record in records:
        slide_number = int(record["source_page_or_slide"])
        slide_path = slide_paths[slide_number]
        asset_type = record.get("asset_type")
        if asset_type == "ppt_slide":
            record["render_path"] = str(slide_path)
            record["sha256"] = _sha256_bytes(slide_path.read_bytes())
            record["delivery_render_sha256"] = record["sha256"]
        elif asset_type in {
            "ppt_picture",
            "ppt_chart_crop",
            "ppt_table_crop",
            "ppt_diagram_crop",
        }:
            bbox = list(record.get("bbox_emu") or [])
            if not valid_bbox_xywh(bbox):
                record["asset_scope"] = "rejected"
                record["deliverable"] = False
                record["render_path"] = None
                record["extraction_error"] = "invalid_object_geometry"
                record["is_decorative"] = True
                continue
            target = crop_dir / f"{record['asset_id']}.png"
            crop_shape_from_slide(
                slide_path,
                target,
                bbox_emu=bbox,
                slide_size_emu=slide_size_emu,
            )
            record["render_path"] = str(target.resolve())
            record["sha256"] = _sha256_bytes(target.read_bytes())
            record["delivery_render_sha256"] = record["sha256"]
            record["delivery_visual_phash"] = _visual_hash(target)
    return deduplicate_assets(records)


def finalize_inventory_records(
    records: list[dict[str, object]], *, source_file_count: int
) -> tuple[list[dict[str, object]], dict[str, int]]:
    records = deduplicate_assets(records)
    for record in records:
        if record.get("asset_type") == "pdf_image_crop":
            original_value = record.get("original_media_path")
            if original_value and Path(str(original_value)).exists():
                original_path = Path(str(original_value)).resolve()
                delivery_hash = _sha256_bytes(original_path.read_bytes())
                record["render_path"] = str(original_path)
                record["sha256"] = delivery_hash
                record["delivery_render_sha256"] = delivery_hash
                record["delivery_visual_phash"] = _visual_hash(original_path)
                record["delivery_basis"] = "original_embedded_media"
        contract = validate_deliverable_asset(record)
        record["contract_errors"] = contract["errors"]
        quality = evaluate_asset_quality(record)
        record.update(quality)
        if "identity_visual_mismatch" in quality["quality_warnings"]:
            if float(record.get("area_ratio") or 0) < 0.02:
                record["asset_role"] = "decoration"
                record["is_decorative"] = True
                record["role_reason"] = ["identity_geometry_but_photo_like_small_asset"]
            else:
                record["asset_role"] = "content_visual"
                record["is_decorative"] = False
                record["role_reason"] = ["identity_geometry_but_substantive_photo"]
            record["role_confidence"] = 0.95
        forbidden_errors = {
            "stamp_or_seal",
            "text_dominant_image",
            "qr_code_or_scanner_mark",
        }
        if forbidden_errors.intersection(quality["quality_errors"]):
            record["asset_role"] = "forbidden"
            record["role_reason"] = sorted(
                set(record.get("role_reason", [])).union(
                    forbidden_errors.intersection(quality["quality_errors"])
                )
            )
            record["role_confidence"] = 1.0
        record["eligible"] = bool(
            contract["valid"]
            and quality["quality_status"] in {"pass", "pass_with_warning"}
            and not record.get("is_decorative")
            and record.get("duplicate_of") is None
        )

    summary = {
        "source_files": source_file_count,
        "assets": len(records),
        "eligible_assets": sum(bool(record["eligible"]) for record in records),
        "quality_rejected": sum(record.get("quality_status") == "reject" for record in records),
        "evidence_a": sum(record.get("evidence_level") == "A" for record in records),
        "evidence_b": sum(record.get("evidence_level") == "B" for record in records),
        "pdf_pages": sum(record["asset_type"] == "pdf_page" for record in records),
        "pdf_image_crops": sum(record["asset_type"] == "pdf_image_crop" for record in records),
        "pdf_chart_crops": sum(record["asset_type"] == "pdf_chart_crop" for record in records),
        "pdf_table_crops": sum(record["asset_type"] == "pdf_table_crop" for record in records),
        "pdf_diagram_crops": sum(record["asset_type"] == "pdf_diagram_crop" for record in records),
        "ppt_slides": sum(record["asset_type"] == "ppt_slide" for record in records),
        "ppt_pictures": sum(record["asset_type"] == "ppt_picture" for record in records),
        "ppt_picture_unreadable": sum(record["asset_type"] == "ppt_picture_unreadable" for record in records),
        "ppt_chart_crops": sum(record["asset_type"] == "ppt_chart_crop" for record in records),
        "ppt_table_crops": sum(record["asset_type"] == "ppt_table_crop" for record in records),
        "ocr_required": sum(bool(record.get("ocr_required")) for record in records),
        "decorative": sum(bool(record.get("is_decorative")) for record in records),
        "duplicates": sum(record.get("duplicate_of") is not None for record in records),
        "identity_candidates": sum(record.get("asset_role") == "identity_candidate" for record in records),
        "forbidden": sum(record.get("asset_role") == "forbidden" for record in records),
    }
    return records, summary


def build_inventory(
    source_files: list[str | Path],
    output_dir: str | Path,
    *,
    render: bool = True,
) -> dict[str, object]:
    destination = Path(output_dir)
    destination.mkdir(parents=True, exist_ok=True)
    records: list[dict[str, object]] = []

    for source_value in source_files:
        source = Path(source_value)
        source_dir = destination / "sources" / _safe_stem(source)
        suffix = source.suffix.lower()
        if suffix == ".pdf":
            records.extend(extract_pdf_pages(source, source_dir, render=render))
        elif suffix == ".pptx":
            ppt_records = extract_pptx_structure(source, source_dir)
            if render:
                slide_dir = source_dir / "slides"
                rendered = render_pptx_slides(source, slide_dir)
                deck = Presentation(source)
                ppt_records = attach_pptx_renders(
                    ppt_records,
                    rendered,
                    slide_size_emu=[deck.slide_width, deck.slide_height],
                    output_dir=source_dir / "shape-crops",
                )
            records.extend(ppt_records)
        else:
            raise ValueError(f"Unsupported source format: {source}")

    records, summary = finalize_inventory_records(
        records, source_file_count=len(source_files)
    )
    payload = {
        "schema_version": "ppt-word-material-inventory-v2",
        "summary": summary,
        "assets": records,
    }
    (destination / "manifest.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return payload


def validate_inventory(payload: dict[str, object]) -> dict[str, object]:
    errors: list[str] = []
    seen_ids: set[str] = set()
    checked = 0
    for record in payload.get("assets", []):
        asset_id = str(record.get("asset_id") or "")
        if not asset_id or asset_id in seen_ids:
            errors.append(f"duplicate_or_empty_asset_id:{asset_id}")
        seen_ids.add(asset_id)
        if not record.get("eligible"):
            continue
        checked += 1
        render_value = record.get("render_path")
        if not render_value:
            errors.append(f"missing_render:{asset_id}")
            continue
        render_path = Path(str(render_value))
        if not render_path.exists():
            errors.append(f"render_not_found:{asset_id}:{render_path}")
            continue
        expected_hash = str(record.get("sha256") or "")
        actual_hash = _sha256_bytes(render_path.read_bytes())
        if expected_hash != actual_hash:
            errors.append(f"sha256_mismatch:{asset_id}")
        if record.get("duplicate_of") is not None:
            errors.append(f"eligible_duplicate:{asset_id}")
    return {"valid": not errors, "checked_eligible": checked, "errors": errors}
